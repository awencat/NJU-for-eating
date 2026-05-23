from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "ppt" / "template-edit" / "source.pptx"
OUT = Path.home() / "Downloads" / "智慧校园餐厅推荐系统汇报-最终版-沿用原模板.pptx"

DARK = "2D3748"
MUTED = "718096"
GREEN = "2ECC71"
LIGHT_GREEN = "F0FFF4"
WHITE = "FFFFFF"


def rgb(hex_color):
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def set_text(shape, text):
    if not hasattr(shape, "text_frame"):
        return
    shape.text_frame.clear()
    p = shape.text_frame.paragraphs[0]
    p.text = text


def set_font(paragraph, size, color=DARK, bold=False):
    for run in paragraph.runs:
        run.font.name = "Noto Sans SC"
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
        run.font.bold = bold


def replace_text(prs, old, new):
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and old in shape.text:
                shape.text = shape.text.replace(old, new)


def remove_slide(prs, index):
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide_id_list.remove(slides[index])


def add_title(slide, title, subtitle):
    title_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.83), Inches(0.69), Inches(11.67), Inches(0.56))
    title_box.fill.background()
    title_box.line.fill.background()
    title_box.text = title
    p = title_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    set_font(p, 32, DARK, True)

    sub_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.83), Inches(1.36), Inches(11.67), Inches(0.33))
    sub_box.fill.background()
    sub_box.line.fill.background()
    sub_box.text = subtitle
    sp = sub_box.text_frame.paragraphs[0]
    sp.alignment = PP_ALIGN.LEFT
    set_font(sp, 16, MUTED, False)


def add_panel(slide, x, y, w, h, title, lines):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb(WHITE)
    panel.line.color.rgb = rgb("E2E8F0")

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.14), Inches(y + 0.16), Inches(0.12), Inches(h - 0.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(GREEN)
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(x + 0.38), Inches(y + 0.2), Inches(w - 0.56), Inches(h - 0.3))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(6)
    set_font(p, 15, DARK, True)
    for line in lines:
        bp = tf.add_paragraph()
        bp.text = line
        bp.level = 0
        bp.space_after = Pt(4)
        set_font(bp, 11.2, DARK, False)
        bp._p.get_or_add_pPr().insert(0, __import__("pptx").oxml.parse_xml(
            '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="•"/>'
        ))


def add_value_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "项目价值", "从餐厅推荐延伸到校园短距离出行决策")

    add_panel(slide, 0.95, 2.05, 3.55, 3.95, "对用户的价值", [
        "把评分、价格、距离、等待时间放在同一界面比较",
        "根据个人预算、菜系偏好和排队接受程度生成推荐",
        "直接查看路线和预计耗时，减少临时决策成本",
    ])
    add_panel(slide, 4.9, 2.05, 3.55, 3.95, "对课程主题的价值", [
        "将位置服务、路径规划和出行时间估计落到校园生活场景",
        "体现智慧交通中“出行目的地选择 + 路径选择”的结合",
        "用轻量系统验证校园短距离出行优化思路",
    ])
    add_panel(slide, 8.85, 2.05, 3.55, 3.95, "对工程实践的价值", [
        "完成前后端分离、接口联调和本地数据管理",
        "处理 WGS-84 与 GCJ-02 坐标系统差异",
        "设计 API 降级逻辑，提高演示和运行稳定性",
    ])


def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "总结与反思", "项目最终完成情况与开发收获")

    add_panel(slide, 0.95, 2.0, 5.55, 4.15, "最终完成内容", [
        "完成餐厅数据管理、搜索筛选、个性化推荐和地图展示",
        "实现步行、骑行、公交三种出行方式的路径规划",
        "完成推荐理由生成、收藏设置和偏好本地持久化",
        "修正坐标系统问题，保证定位、标记和路线显示一致",
    ])
    add_panel(slide, 6.95, 2.0, 5.45, 4.15, "项目收获", [
        "理解了智慧交通出行中位置、路径和用户偏好的关系",
        "掌握了规则型推荐算法在小规模数据场景下的实现方式",
        "认识到地图服务接入、坐标转换和异常降级对系统可用性的影响",
        "完成了从需求分析到系统实现、联调和汇报展示的完整过程",
    ])


def move_slide_to_end(prs, index):
    slide_id_list = prs.slides._sldIdLst
    slide_ids = list(slide_id_list)
    slide_id_list.remove(slide_ids[index])
    slide_id_list.append(slide_ids[index])


def main():
    prs = Presentation(SRC)

    # 补目录第 04 项
    contents = prs.slides[1]
    if len(contents.shapes) > 13 and hasattr(contents.shapes[13], "text_frame"):
        set_text(contents.shapes[13], "项目价值与总结")
        p = contents.shapes[13].text_frame.paragraphs[0]
        set_font(p, 18, DARK, True)

    # 修正中期稿里不适合最终汇报的旧表述
    replace_text(prs, "数据库包含400家南京大学鼓楼校区周边餐厅，信息实时更新。", "数据库包含198家南京大学鼓楼校区周边餐厅，覆盖名称、位置、菜系、价格、评分和等待时间等字段。")
    replace_text(prs, "基于多维度用户画像算法，提供千人千面的个性化餐厅推荐。", "基于位置、预算、菜系和等待时间等条件，生成可解释的个性化餐厅推荐。")
    replace_text(prs, "实时排队状况", "等待时间")
    replace_text(prs, "前端界面截图占位符", "前端界面展示：地图定位、餐厅筛选与推荐列表联动")

    # 保留原来的 Q&A 页，新增两页后将 Q&A 移到最后，避免破坏原模板结构
    qa_index = len(prs.slides) - 1
    add_value_slide(prs)
    add_summary_slide(prs)
    move_slide_to_end(prs, qa_index)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()

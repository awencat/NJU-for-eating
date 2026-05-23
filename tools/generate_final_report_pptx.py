from pathlib import Path
import math
import textwrap

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Circle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "ppt" / "final-report-pptx" / "assets"
OUTPUT = Path.home() / "Downloads" / "智慧校园餐厅推荐系统最终汇报.pptx"

NAVY = "17365D"
BLUE = "2F75B5"
CYAN = "5B9BD5"
LIGHT = "EAF2F8"
PAPER = "F7F9FC"
GRAY = "5B6770"
DARK = "1F2933"
GREEN = "70AD47"
ORANGE = "ED7D31"
RED = "C00000"


def rgb(hex_color):
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def ensure_assets():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def setup_matplotlib():
    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False


def save_architecture(path):
    setup_matplotlib()
    fig, ax = plt.subplots(figsize=(12, 6.6), dpi=180)
    ax.set_facecolor("#F7F9FC")
    fig.patch.set_facecolor("#F7F9FC")
    ax.axis("off")

    layers = [
        ("前端层", "Leaflet 地图 / 搜索筛选 / 收藏设置 / 路线绘制"),
        ("API 层", "Flask REST API / 参数校验 / 统一 JSON 返回"),
        ("核心层", "推荐算法 / 路径规划 / 偏好过滤 / 坐标转换"),
        ("数据层", "SQLite 餐厅库 / 198 家餐厅 / 菜系价格评分等待时间"),
    ]
    y_positions = [0.78, 0.57, 0.36, 0.15]
    colors = ["#DCEAF7", "#E9F3EA", "#FFF2CC", "#FCE4D6"]
    for idx, ((title, desc), y, color) in enumerate(zip(layers, y_positions, colors)):
        box = FancyBboxPatch((0.12, y), 0.76, 0.13, boxstyle="round,pad=0.018,rounding_size=0.025",
                             linewidth=1.6, edgecolor="#17365D", facecolor=color)
        ax.add_patch(box)
        ax.text(0.17, y + 0.083, title, fontsize=17, fontweight="bold", color="#17365D", va="center")
        ax.text(0.33, y + 0.083, desc, fontsize=13, color="#1F2933", va="center")
        if idx < len(layers) - 1:
            ax.add_patch(FancyArrowPatch((0.5, y - 0.01), (0.5, y_positions[idx + 1] + 0.145),
                                         arrowstyle="-|>", mutation_scale=18, linewidth=1.4, color="#5B6770"))
    ax.text(0.5, 0.96, "系统分层架构", ha="center", va="center", fontsize=20, fontweight="bold", color="#17365D")
    ax.text(0.5, 0.035, "用户请求从前端进入 API，核心层完成推荐与路径计算，数据层提供餐厅基础信息",
            ha="center", va="center", fontsize=11.5, color="#5B6770")
    fig.savefig(path, bbox_inches="tight", pad_inches=0.15)
    plt.close(fig)


def save_weight_chart(path):
    setup_matplotlib()
    labels = ["评分", "价格", "等待时间", "500m内加分"]
    values = [50, 30, 20, 10]
    colors = ["#2F75B5", "#70AD47", "#ED7D31", "#A5A5A5"]
    fig, ax = plt.subplots(figsize=(9, 5.2), dpi=180)
    fig.patch.set_facecolor("#F7F9FC")
    ax.set_facecolor("#F7F9FC")
    bars = ax.barh(labels, values, color=colors, height=0.55)
    ax.set_xlim(0, 60)
    ax.set_xlabel("权重或加分幅度", color="#5B6770")
    ax.set_title("推荐得分构成", fontsize=17, fontweight="bold", color="#17365D", pad=16)
    ax.grid(axis="x", alpha=0.22)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.tick_params(axis="y", labelsize=12)
    ax.tick_params(axis="x", colors="#5B6770")
    for bar, val in zip(bars, values):
        ax.text(val + 1, bar.get_y() + bar.get_height() / 2, f"{val}%", va="center", fontsize=11, color="#1F2933")
    ax.text(0, -0.85, "硬约束：超过最大距离或最高预算的餐厅不进入排序",
            fontsize=10.5, color="#5B6770")
    fig.savefig(path, bbox_inches="tight", pad_inches=0.15)
    plt.close(fig)


def save_route_flow(path):
    setup_matplotlib()
    fig, ax = plt.subplots(figsize=(12, 5.8), dpi=180)
    fig.patch.set_facecolor("#F7F9FC")
    ax.set_facecolor("#F7F9FC")
    ax.axis("off")

    nodes = [
        ("用户选择餐厅", "origin / destination / mode"),
        ("坐标转换", "WGS-84 → GCJ-02"),
        ("高德路径 API", "walking / biking / transit"),
        ("解析路线", "distance / duration / polyline"),
        ("前端绘制", "Leaflet polyline"),
    ]
    x_positions = [0.08, 0.27, 0.47, 0.67, 0.86]
    for i, ((title, desc), x) in enumerate(zip(nodes, x_positions)):
        box = FancyBboxPatch((x - 0.075, 0.46), 0.15, 0.18, boxstyle="round,pad=0.015,rounding_size=0.025",
                             linewidth=1.5, edgecolor="#17365D", facecolor="#EAF2F8")
        ax.add_patch(box)
        ax.text(x, 0.575, title, ha="center", va="center", fontsize=12.5, fontweight="bold", color="#17365D")
        ax.text(x, 0.51, desc, ha="center", va="center", fontsize=9.5, color="#5B6770")
        if i < len(nodes) - 1:
            ax.add_patch(FancyArrowPatch((x + 0.085, 0.55), (x_positions[i + 1] - 0.085, 0.55),
                                         arrowstyle="-|>", mutation_scale=16, linewidth=1.4, color="#5B6770"))
    fallback = FancyBboxPatch((0.39, 0.16), 0.33, 0.15, boxstyle="round,pad=0.018,rounding_size=0.025",
                              linewidth=1.4, edgecolor="#ED7D31", facecolor="#FFF2CC")
    ax.add_patch(fallback)
    ax.text(0.555, 0.245, "异常降级：API 不可用时使用 Haversine 距离 + 30段插值路径",
            ha="center", va="center", fontsize=11.5, fontweight="bold", color="#7F3F00")
    ax.add_patch(FancyArrowPatch((0.47, 0.45), (0.47, 0.31), arrowstyle="-|>", mutation_scale=15,
                                 linewidth=1.2, color="#ED7D31", linestyle="--"))
    ax.add_patch(FancyArrowPatch((0.72, 0.235), (0.86, 0.45), arrowstyle="-|>", mutation_scale=15,
                                 linewidth=1.2, color="#ED7D31", linestyle="--"))
    ax.text(0.5, 0.88, "多模式路径规划与降级机制", ha="center", va="center",
            fontsize=19, fontweight="bold", color="#17365D")
    fig.savefig(path, bbox_inches="tight", pad_inches=0.15)
    plt.close(fig)


def save_coordinate_flow(path):
    setup_matplotlib()
    fig, ax = plt.subplots(figsize=(10, 5.5), dpi=180)
    fig.patch.set_facecolor("#F7F9FC")
    ax.axis("off")
    ax.set_facecolor("#F7F9FC")
    items = [
        (0.18, 0.60, "浏览器定位", "WGS-84", "#DCEAF7"),
        (0.50, 0.60, "统一转换", "wgs84_to_gcj02()", "#E9F3EA"),
        (0.82, 0.60, "地图与餐厅点", "GCJ-02", "#FCE4D6"),
        (0.50, 0.24, "距离计算", "mixed_system_distance()", "#FFF2CC"),
    ]
    for x, y, title, desc, color in items:
        ax.add_patch(FancyBboxPatch((x - 0.13, y - 0.08), 0.26, 0.16,
                                    boxstyle="round,pad=0.018,rounding_size=0.025",
                                    linewidth=1.5, edgecolor="#17365D", facecolor=color))
        ax.text(x, y + 0.025, title, ha="center", va="center", fontsize=13, fontweight="bold", color="#17365D")
        ax.text(x, y - 0.035, desc, ha="center", va="center", fontsize=10.5, color="#5B6770")
    ax.add_patch(FancyArrowPatch((0.31, 0.60), (0.37, 0.60), arrowstyle="-|>", mutation_scale=16, color="#5B6770"))
    ax.add_patch(FancyArrowPatch((0.63, 0.60), (0.69, 0.60), arrowstyle="-|>", mutation_scale=16, color="#5B6770"))
    ax.add_patch(FancyArrowPatch((0.50, 0.51), (0.50, 0.34), arrowstyle="-|>", mutation_scale=16, color="#5B6770"))
    ax.add_patch(Circle((0.50, 0.60), 0.17, fill=False, linewidth=1.2, edgecolor="#70AD47", linestyle="--"))
    ax.text(0.5, 0.90, "坐标系统一致性处理", ha="center", fontsize=18, fontweight="bold", color="#17365D")
    ax.text(0.5, 0.08, "目的：避免用户位置、餐厅标记和路线绘制出现偏移或重复转换",
            ha="center", fontsize=10.5, color="#5B6770")
    fig.savefig(path, bbox_inches="tight", pad_inches=0.15)
    plt.close(fig)


def save_ui_mock(path):
    setup_matplotlib()
    fig, ax = plt.subplots(figsize=(11, 6.2), dpi=180)
    fig.patch.set_facecolor("#F7F9FC")
    ax.set_facecolor("#F7F9FC")
    ax.axis("off")
    ax.add_patch(FancyBboxPatch((0.05, 0.08), 0.90, 0.82, boxstyle="round,pad=0.015,rounding_size=0.025",
                                linewidth=1.5, edgecolor="#17365D", facecolor="#FFFFFF"))
    ax.add_patch(FancyBboxPatch((0.07, 0.12), 0.56, 0.72, boxstyle="round,pad=0.01,rounding_size=0.02",
                                linewidth=1, edgecolor="#B7C9D6", facecolor="#DCEAF7"))
    for i in range(7):
        ax.plot([0.09, 0.61], [0.18 + i * 0.09, 0.23 + i * 0.07], color="#B7C9D6", linewidth=1)
    for x, y, c in [(0.20, 0.58, "#2F75B5"), (0.35, 0.46, "#70AD47"), (0.50, 0.65, "#ED7D31"), (0.42, 0.30, "#A5A5A5")]:
        ax.add_patch(Circle((x, y), 0.018, color=c))
    ax.add_patch(FancyArrowPatch((0.20, 0.58), (0.35, 0.46), arrowstyle="-", linewidth=3, color="#2F75B5"))
    ax.add_patch(FancyArrowPatch((0.35, 0.46), (0.50, 0.65), arrowstyle="-", linewidth=3, color="#2F75B5"))

    ax.add_patch(FancyBboxPatch((0.67, 0.66), 0.22, 0.14, boxstyle="round,pad=0.012,rounding_size=0.018",
                                linewidth=1.2, edgecolor="#D9E2EC", facecolor="#F7F9FC"))
    ax.text(0.69, 0.75, "推荐结果", fontsize=12, fontweight="bold", color="#17365D")
    ax.text(0.69, 0.70, "评分 4.7  人均 ¥32  距离 420m", fontsize=9.5, color="#5B6770")
    ax.add_patch(FancyBboxPatch((0.67, 0.47), 0.22, 0.14, boxstyle="round,pad=0.012,rounding_size=0.018",
                                linewidth=1.2, edgecolor="#D9E2EC", facecolor="#F7F9FC"))
    ax.text(0.69, 0.56, "筛选条件", fontsize=12, fontweight="bold", color="#17365D")
    ax.text(0.69, 0.51, "预算 / 菜系 / 距离 / 等待", fontsize=9.5, color="#5B6770")
    ax.add_patch(FancyBboxPatch((0.67, 0.28), 0.22, 0.14, boxstyle="round,pad=0.012,rounding_size=0.018",
                                linewidth=1.2, edgecolor="#D9E2EC", facecolor="#F7F9FC"))
    ax.text(0.69, 0.37, "路线规划", fontsize=12, fontweight="bold", color="#17365D")
    ax.text(0.69, 0.32, "步行 / 骑行 / 公交", fontsize=9.5, color="#5B6770")
    ax.text(0.5, 0.95, "前端交互界面示意", ha="center", fontsize=18, fontweight="bold", color="#17365D")
    fig.savefig(path, bbox_inches="tight", pad_inches=0.15)
    plt.close(fig)


def add_bg(slide, color=PAPER):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(color)


def add_footer(slide, page, title="智慧校园餐厅推荐系统"):
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(6.2), Inches(0.22))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(8.5)
    p.font.color.rgb = rgb(GRAY)
    p.alignment = PP_ALIGN.LEFT
    tx2 = slide.shapes.add_textbox(Inches(12.0), Inches(7.05), Inches(0.75), Inches(0.22))
    p2 = tx2.text_frame.paragraphs[0]
    p2.text = f"{page:02d}"
    p2.font.name = "Arial"
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = rgb(GRAY)
    p2.alignment = PP_ALIGN.RIGHT


def add_title(slide, title, subtitle=None, page=None, dark=False):
    color = "FFFFFF" if dark else NAVY
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.42), Inches(9.7), Inches(0.62))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = rgb(color)
    p.alignment = PP_ALIGN.LEFT
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.62), Inches(1.02), Inches(8.5), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Microsoft YaHei"
        sp.font.size = Pt(10.5)
        sp.font.color.rgb = rgb("D9E2EC" if dark else GRAY)
    if page is not None:
        add_footer(slide, page)


def add_label(slide, text, x, y, w, color=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.22))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(7)
    p.font.bold = True
    p.font.color.rgb = rgb("FFFFFF")
    p.alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("FFFFFF")
    shape.line.color.rgb = rgb("D9E2EC")
    shape.line.width = Pt(1)
    add_label(slide, title, x + 0.18, y + 0.18, min(w - 0.36, 1.3), accent)
    tx = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.58), Inches(w - 0.44), Inches(h - 0.76))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = body
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(13)
    p.font.color.rgb = rgb(DARK)
    p.line_spacing = 1.15
    return shape


def add_filled_shape(slide, shape_type, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    if line_color:
        shape.line.color.rgb = rgb(line_color)
    else:
        shape.line.fill.background()
    return shape


def add_bullets(slide, items, x, y, w, h, size=14, color=DARK, gap=0.04):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(gap * 72)
        p.level = 0
        p._p.get_or_add_pPr().insert(0, __import__("pptx").oxml.parse_xml(
            '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="•"/>'
        ))
    return tx


def add_metric(slide, x, y, number, label, note=None, color=BLUE):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.2), Inches(0.75))
    p = tx.text_frame.paragraphs[0]
    p.text = number
    p.font.name = "Arial"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = rgb(color)
    p.alignment = PP_ALIGN.CENTER
    lab = slide.shapes.add_textbox(Inches(x), Inches(y + 0.72), Inches(2.2), Inches(0.35))
    lp = lab.text_frame.paragraphs[0]
    lp.text = label
    lp.font.name = "Microsoft YaHei"
    lp.font.size = Pt(11)
    lp.font.bold = True
    lp.font.color.rgb = rgb(DARK)
    lp.alignment = PP_ALIGN.CENTER
    if note:
        nt = slide.shapes.add_textbox(Inches(x), Inches(y + 1.05), Inches(2.2), Inches(0.32))
        np = nt.text_frame.paragraphs[0]
        np.text = note
        np.font.name = "Microsoft YaHei"
        np.font.size = Pt(8)
        np.font.color.rgb = rgb(GRAY)
        np.alignment = PP_ALIGN.CENTER


def add_code_box(slide, code, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("17212B")
    shape.line.color.rgb = rgb("2E4052")
    tx = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.16), Inches(w - 0.36), Inches(h - 0.25))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = "Consolas"
    p.font.size = Pt(10)
    p.font.color.rgb = rgb("EAF2F8")
    return shape


def add_section_badge(slide, text):
    add_label(slide, text, 0.62, 0.22, 1.6, NAVY)


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(NAVY)
    bg.line.fill.background()
    add_label(slide, "智慧交通出行 · 课程最终汇报", 0.72, 0.55, 2.4, CYAN)
    tx = slide.shapes.add_textbox(Inches(0.72), Inches(1.65), Inches(10.4), Inches(1.5))
    p = tx.text_frame.paragraphs[0]
    p.text = "智慧校园餐厅推荐系统"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = rgb("FFFFFF")
    sub = slide.shapes.add_textbox(Inches(0.76), Inches(3.1), Inches(9.8), Inches(0.65))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "基于位置、偏好与路径成本的校园餐饮出行优化"
    sp.font.name = "Microsoft YaHei"
    sp.font.size = Pt(19)
    sp.font.color.rgb = rgb("D9E2EC")
    for x, y, r, c in [(10.8, 1.1, 0.26, CYAN), (11.55, 1.9, 0.18, GREEN), (10.55, 3.0, 0.14, ORANGE), (11.35, 3.65, 0.22, "FFFFFF")]:
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(r), Inches(r))
        circ.fill.solid()
        circ.fill.fore_color.rgb = rgb(c)
        circ.line.fill.background()
    meta = slide.shapes.add_textbox(Inches(0.76), Inches(5.4), Inches(7.8), Inches(0.8))
    mp = meta.text_frame.paragraphs[0]
    mp.text = "小组成员：石城玮、李郅涵    日期：2026年5月23日"
    mp.font.name = "Microsoft YaHei"
    mp.font.size = Pt(13)
    mp.font.color.rgb = rgb("D9E2EC")
    add_footer(slide, 1, "NJU-for-eating")


def build_presentation():
    ensure_assets()
    arch = ASSET_DIR / "architecture.png"
    weights = ASSET_DIR / "weights.png"
    route = ASSET_DIR / "route-flow.png"
    coord = ASSET_DIR / "coordinate-flow.png"
    ui = ASSET_DIR / "ui-mock.png"
    save_architecture(arch)
    save_weight_chart(weights)
    save_route_flow(route)
    save_coordinate_flow(coord)
    save_ui_mock(ui)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)

    # 2 目录
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "汇报目录", "从需求、系统实现到结果与改进", 2)
    items = [
        ("01", "项目背景与目标", "说明校园餐饮出行问题和系统目标"),
        ("02", "系统设计与数据基础", "介绍架构、数据字段和功能闭环"),
        ("03", "核心算法与路径规划", "说明推荐评分、坐标转换和路线降级"),
        ("04", "最终成果与不足", "总结已完成内容、问题和后续方向"),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = 1.65 + i * 1.12
        add_metric(slide, 0.95, y - 0.06, num, "", color=BLUE)
        add_card(slide, 2.4, y, 8.8, 0.82, title, desc, [BLUE, GREEN, ORANGE, NAVY][i])

    # 3 背景
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "项目背景与痛点分析", "校园用餐决策同时包含餐厅选择和出行成本", 3)
    add_card(slide, 0.75, 1.62, 3.75, 2.0, "选择困难", "餐厅数量多，评分、人均价格、等待时间等信息分散，学生需要反复比较。", BLUE)
    add_card(slide, 4.8, 1.62, 3.75, 2.0, "出行规划复杂", "步行、骑行、公交适用场景不同，仅按距离或评分排序不能反映真实到达成本。", GREEN)
    add_card(slide, 8.85, 1.62, 3.75, 2.0, "个性化不足", "预算、菜系偏好、是否接受排队会影响选择，需要让筛选条件进入推荐流程。", ORANGE)
    add_bullets(slide, [
        "项目目标：围绕“现在去哪吃更合适”建立完整决策流程。",
        "课程关联：把位置服务、路径规划和出行时间估计引入校园餐饮场景。",
        "最终产出：可运行的前后端系统，而不是单独的静态推荐列表。"
    ], 1.05, 4.35, 11.1, 1.45, 15)

    # 4 数据与范围
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "数据基础与服务范围", "以南京大学鼓楼校区周边餐厅为主要对象", 4)
    for x, num, label, note, color in [
        (0.9, "198", "餐厅数据", "当前项目数据规模", BLUE),
        (3.65, "10+", "菜系类型", "快餐、川菜、西餐、咖啡等", GREEN),
        (6.4, "10-100", "价格区间", "人均价格，单位为元", ORANGE),
        (9.15, "3", "出行方式", "步行、骑行、公交", NAVY),
    ]:
        add_metric(slide, x, 1.48, num, label, note, color)
    add_card(slide, 0.9, 3.55, 5.65, 1.65, "餐厅字段", "名称、地址、坐标、菜系、人均价格、评分、等待时间、营业信息和标签。", BLUE)
    add_card(slide, 6.85, 3.55, 5.65, 1.65, "用户输入", "当前位置、最高预算、最大距离、菜系偏好、是否接受排队等条件。", GREEN)
    add_bullets(slide, [
        "最终汇报统一采用当前项目数据口径：餐厅样本数为 198 家。",
        "等待时间目前是静态字段，后续可接入实时客流或食堂排队数据。"
    ], 1.0, 5.72, 11.5, 0.9, 13)

    # 5 架构
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "系统架构设计", "前端交互、API 服务、核心算法、数据访问分层实现", 5)
    slide.shapes.add_picture(str(arch), Inches(1.05), Inches(1.42), width=Inches(11.3))

    # 6 功能闭环
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "最终功能闭环", "从定位到推荐，再到路线展示和偏好持久化", 6)
    steps = [
        ("定位", "获取当前位置并显示精度范围"),
        ("推荐", "结合预算、菜系、等待时间和距离排序"),
        ("筛选", "按价格、评分、距离、菜系等条件调整结果"),
        ("导航", "展示步行、骑行、公交路径与耗时"),
        ("保存", "收藏餐厅和偏好设置写入 localStorage"),
    ]
    for i, (title, desc) in enumerate(steps):
        x = 0.82 + i * 2.45
        add_card(slide, x, 2.0, 2.05, 2.05, title, desc, [BLUE, GREEN, ORANGE, NAVY, CYAN][i])
        if i < len(steps) - 1:
            add_filled_shape(slide, MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, x + 1.93, 2.72, 0.45, 0.34, "B7C9D6")
    add_bullets(slide, [
        "前端模块：app、map、api、search、filter、favorites、settings。",
        "后端接口：/api/recommend、/api/route、/api/restaurants、/api/health。"
    ], 1.0, 5.12, 11.3, 1.0, 14)

    # 7 推荐功能
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "核心功能一：智能餐厅推荐", "硬约束过滤 + 加权评分排序 + 推荐理由生成", 7)
    add_bullets(slide, [
        "先过滤超过最大距离、最高预算和不可接受等待时间的餐厅。",
        "对剩余餐厅计算评分、价格、等待时间三个维度的归一化得分。",
        "500 米以内餐厅增加距离奖励，贴合校园短距离用餐习惯。",
        "接口默认返回得分最高的前 10 家餐厅，并给出可理解的推荐理由。"
    ], 0.9, 1.65, 5.7, 3.3, 14)
    slide.shapes.add_picture(str(weights), Inches(7.0), Inches(1.5), width=Inches(5.2))
    add_code_box(slide, "score = 0.5 * rating\n      + 0.3 * price\n      + 0.2 * wait\n      + nearby_bonus", 0.95, 5.2, 5.3, 1.0)

    # 8 算法实现
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "推荐算法实现细节", "使用可解释规则，适合当前数据规模和课程演示场景", 8)
    add_code_box(slide, textwrap.dedent("""\
        if distance > max_distance:
            return -inf, reasons
        if restaurant["price"] > max_price:
            return -inf, reasons

        norm_rating = restaurant["rating"] / 5.0
        norm_price = 1 - normalize(price, 0, max_price)
        norm_wait = 1 - normalize(wait_time, 0, 60)
    """), 0.8, 1.45, 6.1, 2.55)
    add_card(slide, 7.35, 1.45, 4.95, 1.25, "为什么不用复杂模型", "当前数据量和反馈数据不足，规则模型更可控，也便于解释和调参。", BLUE)
    add_card(slide, 7.35, 3.0, 4.95, 1.25, "可解释性", "返回“距离很近”“评分较高”“价格实惠”“排队时间短”等原因。", GREEN)
    add_card(slide, 7.35, 4.55, 4.95, 1.25, "可扩展性", "后续可加入真实拥挤度、用户历史偏好和时间段权重。", ORANGE)

    # 9 路径规划
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "核心功能二：多模式路径规划", "把推荐结果转化为可到达方案", 9)
    slide.shapes.add_picture(str(route), Inches(0.85), Inches(1.38), width=Inches(11.7))
    add_bullets(slide, [
        "高德 API 可用时返回真实路线；不可用时仍返回直线路径和估算耗时。",
        "路线结果包含 distance、duration、polyline、mode、provider，前端可直接绘制。"
    ], 1.15, 6.17, 11.0, 0.75, 12)

    # 10 坐标
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "关键技术问题：坐标系统一致性", "解决浏览器定位与高德地图服务之间的坐标差异", 10)
    slide.shapes.add_picture(str(coord), Inches(1.05), Inches(1.35), width=Inches(6.4))
    add_bullets(slide, [
        "浏览器定位通常返回 WGS-84 坐标，高德地图瓦片和路径接口使用 GCJ-02。",
        "后端距离计算先转换用户坐标，再与餐厅坐标计算 Haversine 距离。",
        "前端显示用户位置和移动地图视野时同步转换，避免标记偏移。",
        "餐厅坐标已按高德体系使用，避免对餐厅点重复转换。"
    ], 8.0, 1.7, 4.4, 3.7, 13)
    add_code_box(slide, "user_gcj02 = wgs84_to_gcj02(user_lat, user_lng)\ndistance = haversine_distance(user_gcj02, poi_gcj02)", 8.0, 5.45, 4.4, 0.85)

    # 11 前端
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "前端交互与界面组织", "围绕地图完成推荐、筛选和路径展示", 11)
    slide.shapes.add_picture(str(ui), Inches(0.85), Inches(1.3), width=Inches(6.5))
    add_card(slide, 7.65, 1.45, 4.85, 1.05, "地图交互", "使用 Leaflet 加载高德瓦片，显示用户位置、餐厅标记和路线。", BLUE)
    add_card(slide, 7.65, 2.75, 4.85, 1.05, "搜索筛选", "支持餐厅名称、菜系、地址和标签检索，结合防抖降低重复操作。", GREEN)
    add_card(slide, 7.65, 4.05, 4.85, 1.05, "本地持久化", "偏好设置和收藏保存在 localStorage，刷新页面后仍可保留。", ORANGE)
    add_card(slide, 7.65, 5.35, 4.85, 1.05, "性能处理", "地图启用 Canvas 渲染，并限制标记数量，保证基本流畅度。", NAVY)

    # 12 后端
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "后端接口与工程化处理", "面向前端稳定调用的 Flask 服务", 12)
    add_card(slide, 0.8, 1.45, 3.85, 1.45, "接口设计", "推荐、路径、餐厅、健康检查接口均使用 JSON 格式返回。", BLUE)
    add_card(slide, 4.95, 1.45, 3.85, 1.45, "参数校验", "对经纬度、模式、筛选条件进行校验，减少异常输入影响。", GREEN)
    add_card(slide, 9.1, 1.45, 3.45, 1.45, "跨域配置", "支持前后端分离运行，方便本地调试和课堂演示。", ORANGE)
    add_code_box(slide, textwrap.dedent("""\
        {
          "code": 200,
          "status": "success",
          "message": "推荐成功",
          "data": [...],
          "timestamp": 1779...
        }
    """), 0.95, 3.55, 5.45, 2.2)
    add_bullets(slide, [
        "配置项集中在 config.py 中管理，包括 API Key、端口、数据库路径和 CORS 源。",
        "路径规划服务失败时降级到简单路径，提升演示环境下的可用性。",
        "SQLite 适合当前课程项目规模，后续可迁移到服务端数据库。"
    ], 7.0, 3.7, 5.2, 1.85, 13)

    # 13 成果
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "最终成果展示", "项目已经形成可运行的校园餐饮出行推荐系统", 13)
    for x, num, label, note, color in [
        (0.9, "4", "后端核心模块", "推荐、路径、筛选、数据访问", BLUE),
        (3.8, "7", "前端 JS 模块", "地图、搜索、筛选、收藏等", GREEN),
        (6.7, "3", "出行模式", "步行、骑行、公交", ORANGE),
        (9.6, "1", "完整流程", "定位到推荐再到路径", NAVY),
    ]:
        add_metric(slide, x, 1.45, num, label, note, color)
    add_bullets(slide, [
        "完成餐厅数据导入、推荐排序、地图展示、路线绘制和偏好保存。",
        "修复坐标系统混用问题，统一处理 WGS-84 与 GCJ-02 的转换。",
        "系统支持无高德 Key 的降级运行，保证基本演示流程可完成。",
        "相比中期版本，最终版更强调路径成本和工程可用性。"
    ], 1.0, 4.0, 11.2, 1.9, 15)

    # 14 挑战
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "困难、限制与后续优化", "明确当前系统边界，给出可落地的改进方向", 14)
    add_card(slide, 0.85, 1.45, 3.75, 2.0, "定位精度", "室内定位和网络定位存在波动，后续可加入手动校准和常用起点选择。", BLUE)
    add_card(slide, 4.85, 1.45, 3.75, 2.0, "实时数据", "等待时间目前不是实时采集，后续可接入排队估计或用户反馈。", ORANGE)
    add_card(slide, 8.85, 1.45, 3.75, 2.0, "路径质量", "API 配额和网络状态会影响真实路径，后续可缓存常用路线。", GREEN)
    add_bullets(slide, [
        "推荐算法可以加入时间段因素，例如午餐高峰、晚餐高峰和夜间营业状态。",
        "可建立用户反馈闭环，用收藏、评分和点击记录动态调整权重。",
        "可从餐厅推荐扩展到校园出行服务，例如校内活动、打印点、快递点导航。"
    ], 1.0, 4.25, 11.2, 1.5, 14)

    # 15 结尾
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_label(slide, "Conclusion", 0.75, 0.65, 1.4, CYAN)
    tx = slide.shapes.add_textbox(Inches(0.78), Inches(1.75), Inches(10.8), Inches(1.35))
    p = tx.text_frame.paragraphs[0]
    p.text = "从“找餐厅”到“规划一次用餐出行”"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = rgb("FFFFFF")
    sub = slide.shapes.add_textbox(Inches(0.82), Inches(3.35), Inches(10.2), Inches(1.15))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "本项目将餐厅推荐、位置服务和路径规划整合到同一流程中，完成了一个面向校园场景的智慧交通出行应用原型。"
    sp.font.name = "Microsoft YaHei"
    sp.font.size = Pt(18)
    sp.font.color.rgb = rgb("D9E2EC")
    qa = slide.shapes.add_textbox(Inches(0.82), Inches(5.35), Inches(3.0), Inches(0.6))
    qp = qa.text_frame.paragraphs[0]
    qp.text = "Q & A"
    qp.font.name = "Arial"
    qp.font.size = Pt(24)
    qp.font.bold = True
    qp.font.color.rgb = rgb("FFFFFF")
    add_footer(slide, 15, "最终汇报")

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    out = build_presentation()
    print(out)
